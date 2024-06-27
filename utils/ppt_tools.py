# -*- coding: utf-8 -*-
"""
@Time    : 2024/3/28 16:24
@Author  : minglang.wu
@Email   : minglang.wu@tenclass.com
@File    : ppt_tools.py
@Desc    :
"""
import os
import time
import uuid
import copy
from pptx.enum.shapes import MSO_SHAPE_TYPE


def recreate_slide_by_win32(ppt_path: str, save_path: str, indexs: list):
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        raise ImportError("Please install pywin32 first, pip install pywin32")

    ppt_instance, prs = None, None
    try:
        # Avoid the error of "CoInitialize has not been called"
        pythoncom.CoInitialize()
        ppt_instance = win32com.client.Dispatch("PowerPoint.Application")
        # open the powerpoint presentation headless in background
        read_only = True
        has_title = False
        window = False
        try:
            prs = ppt_instance.Presentations.open(
                os.path.join(os.getcwd(), ppt_path), read_only, has_title, window
            )
            # Record the count of template slides
            count = prs.Slides.Count
            print("当前页面数量：", prs.Slides.Count, count)
            # Copy and paste the target slides
            for index in indexs:
                prs.Slides(index + 1).Copy()
                prs.Slides.Paste()  # Index=insert_index, 不填则默认在最后插入
                print("当前页面数量：", prs.Slides.Count, count)
                time.sleep(0.3)  # 防止复制过快导致丢失

            # Delete template slides
            for _ in range(count):
                prs.Slides(1).Delete()
            prs.SaveAs(os.path.join(os.getcwd(), save_path))
        finally:
            if prs:
                prs.Close()
    finally:
        if ppt_instance:
            ppt_instance.Quit()
            del ppt_instance
        pythoncom.CoUninitialize()


def duplicate_slide(pres, slide_index):
    """
    duplicate side base on one slide

    refs from stackoverflow:
        https://stackoverflow.com/a/73954830
        https://stackoverflow.com/a/56074651/20159015
        https://stackoverflow.com/a/62921848/20159015
        https://stackoverflow.com/questions/50866634/how-to-copy-a-slide-with-python-pptx
        https://stackoverflow.com/questions/62864082/how-to-copy-a-slide-with-images-using-python-pptx
    refs from python-pptx github issue:
        https://github.com/scanny/python-pptx/issues/132
        https://github.com/scanny/python-pptx/issues/238
    """

    def get_uuid():
        return str(uuid.uuid4())

    slide_to_copy = pres.slides[slide_index]

    new_slide = pres.slides.add_slide(slide_to_copy.slide_layout)

    # create images dict
    imgDict = {}

    for shp in slide_to_copy.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_name = get_uuid()
            # save image
            with open(img_name + ".jpg", "wb") as f:
                f.write(shp.image.blob)

            # add image to dict
            imgDict[img_name + ".jpg"] = [shp.left, shp.top, shp.width, shp.height]
        else:
            # create copy of elem
            el = shp.element
            newel = copy.deepcopy(el)

            # add elem to shape tree
            new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

    # things added first will be covered by things added last => since I want pictures to be in foreground, I will add them after others elements
    # you can change this if you want
    # add pictures
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide  # this returns slide so you can instantly work with it when it is pasted in presentation
