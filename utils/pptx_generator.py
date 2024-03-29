import json
import random
import re
import time
import logging
import pptx
from pptx import Presentation

from utils.llm import LLM
from utils.ppt_tools import recreate_slide_by_win32
from utils.prompter import PromptLibrarian

logger = logging.getLogger(__name__)


class PptxGenerator:
    PPT_PARAM_PATTERN = r'\{(.*?)\}'
    MD_CODE_JSON_PATTERN = r'```json.*?\n(.*?)```'

    def __init__(self, llm: LLM, save_path: str, template_path: str = None):
        self.llm = llm
        self.save_path = save_path
        self.has_template = False
        if template_path:
            self.has_template = True
            self.template_path = template_path
            self.template_ppt = Presentation(template_path)
            self.template_params = self._extract_params_from_template()

    def _extract_params_from_template(self):
        """生成PPT文件"""
        # 注意提取参数时，需要把模板PPT中的<组合>都解锁，不然可能存在找不到文本框的情况
        start_slide_idx = 0
        catalogue_slide_idx = 1
        title_slide_idx = 2
        content_slide_idxs = [3, 4, 5, 6, 7, 8]
        end_slide_idx = 9
        template_params = {
            "first_slide": {"nos": [start_slide_idx], "params": []},
            "catalogue_slide": {"nos": [catalogue_slide_idx], "params": []},
            "title_slide": {"nos": [title_slide_idx], "params": []},
            "content_slide": {"nos": content_slide_idxs, "params": []},
            "end_slide": {"nos": [end_slide_idx], "params": []}
        }

        # PPT中同一页的{params}定义必须不同，避免混淆，不同页面不做要求
        for slide_name, slide_info in template_params.items():
            nos = slide_info["nos"]
            for n in nos:
                slide = self.template_ppt.slides[n]
                temp_params = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                matches = [match.group(1) for match in re.finditer(self.PPT_PARAM_PATTERN, run.text)]
                                temp_params.extend(matches)
                slide_info["params"].append(temp_params)
        return template_params

    def llm_generate_online_content(self, topic: str):
        """根据主题生成大纲"""
        output_format = json.dumps(
            {
                "topic": "str",
                "pages": [
                    {
                        "title": "str",
                        "pages": [
                            {"sub_title": "str", "desc": "str", "content": "str"}
                        ]
                    }
                ]
            }, ensure_ascii=False
        )
        # remove space to save token
        output_format = output_format.replace(" ", "")
        prompt = PromptLibrarian.read(path="ppt.generate_content.v1").format(topic=topic,
                                                                             language="中文",
                                                                             output_format=output_format)
        messages = [
            {"role": "system", "content": "你是个全能助手"},
            {"role": "user", "content": prompt}
        ]
        c = self.llm.chat_in_all(messages)
        if c[-1] != "}":
            logger.warning("[Continuous] Reply not end, go on ...")
            messages.append({"role": "assistant", "content": c})
            messages.append({"role": "user", "content": "继续"})
            c += self.llm.chat_in_all(messages)
        return c

    def _llm_generate_content_slide_in_single(self, prompt: str, temperature: float, tp: dict):
        is_match = True
        try_count = 0
        while try_count <= 3:
            try_count += 1
            ctx = self.llm.chat_once(prompt=prompt, temperature=temperature)
            m = re.findall(self.MD_CODE_JSON_PATTERN, ctx, re.DOTALL)
            if m: ctx = m[0].replace("'", '"')
            # try to load json
            try:
                ctx = json.loads(ctx)
            except Exception as e:
                logger.warning(f"ctx json.loads error: \n{ctx}")
                time.sleep(0.8 * try_count)
                continue
            # try to match params
            gcs = [gc for gc in ctx.keys()]
            for tk in tp.keys():
                if tk not in gcs:
                    is_match = False
                    break
            if is_match:
                logger.info(f"try generated count <{try_count}>, ctx: \n{ctx}")
                return ctx
            time.sleep(0.8 * try_count)
        return None

    def llm_generate_content_slide_content(self, topic: str, online_content: str):
        """根据大纲生成完整内容"""
        logger.info(f"online_content: \n{online_content}")
        online_content = json.loads(online_content)
        current_online_content = online_content["pages"]
        content_slide = self.template_params["content_slide"]

        # 新增标题编号、子标题编号
        for idx, c in enumerate(current_online_content):
            c["no"] = idx + 1
        for c in current_online_content:
            for idx, s in enumerate(c["pages"]):
                s["sub_no"] = idx + 1

        title_count = len(current_online_content)
        resorted_no_idxs = random.sample(range(len(content_slide["nos"])), k=title_count)
        current_template_resort_nos = [content_slide["nos"][idx] for idx in resorted_no_idxs]
        logger.info(f"content_slide_params: {content_slide['params']}")
        logger.info(f"resorted_no_idxs: {resorted_no_idxs}")
        logger.info(f"current_template_resort_nos: {current_template_resort_nos}")

        current_template_params = [
            {k: "" for k in content_slide["params"][idx]}
            for idx in resorted_no_idxs
        ]

        for oc, tp in zip(current_online_content, current_template_params):
            title = oc["title"]
            prompt = f"""# Info
## OnlineJson
```{oc}```
## TemplateParamsJson
```{tp}```
# Tasks
严格参照[Info.TemplateParamsJson]，基于《{topic}》中的`{title}`标题的内容，对应填充[Info.OnlineJson]，最后按照markdown的json格式输出。
注意：json的key值严格对应[Info.TemplateParamsJson]，key对应的值不能存在列表或字典。
------
output:"""
            ctx = self._llm_generate_content_slide_in_single(prompt, 0.6, tp)
            if ctx:
                # 严格按照template参数匹配赋值
                for tk in tp.keys():
                    tp[tk] = ctx.get(tk, "")
                time.sleep(2)
            else:
                logger.exception(f"failed to generate content for title: {title}. Skip it!")

        data = {
            "titles_param": {f'title_{i + 1}': c["title"] for i, c in enumerate(current_online_content)},
            "contents_param": current_template_params,
            "nos": current_template_resort_nos
        }
        return data

    def generate_ppt(self, meta_info: dict, generation_content: dict):
        """
        generate ppt based on content which is generated by llm
        :param meta_info:
        :param generation_content:
        :return:
        """
        # 1. 根据模板新开ppt
        logger.info(f"meta_info: {meta_info}")
        logger.info(f"generation_content: {generation_content}")

        titles_param = generation_content["titles_param"]
        contents_param = generation_content["contents_param"]
        nos = generation_content["nos"]

        all_params = contents_param
        # 插入首页内容和slide No
        all_params.insert(0, meta_info)
        nos.insert(0, self.template_params["first_slide"]["nos"][0])
        # 插入目录页
        all_params.insert(1, titles_param)
        nos.insert(1, self.template_params["catalogue_slide"]["nos"][0])
        # 插入结束页(在最后)
        all_params.append({})
        nos.append(self.template_params["end_slide"]["nos"][0])

        # 2. 根据重排的slide No重新生成ppt
        logger.info(f"nos: {nos}")
        recreate_slide_by_win32(self.template_path, self.save_path, indexs=nos)

        # 3. 在新ppt上，新增页并填充内容
        new_ppt = pptx.Presentation(self.save_path)
        for idx, p_dict in enumerate(all_params):
            for shape in new_ppt.slides[idx].shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for match in re.findall(pattern=self.PPT_PARAM_PATTERN, string=run.text):
                                m_str = "{" + match + "}"
                                m_key = match
                                run.text = run.text.replace(m_str, str(p_dict.get(m_key, '')))

        # 4. 保存PPT
        new_ppt.save(self.save_path)

    def generate(self, meta_info: dict):
        """根据模板生成PPT"""
        online_content = self.llm_generate_online_content(meta_info["topic"])
        generation_content = self.llm_generate_content_slide_content(meta_info["topic"], online_content)
        self.generate_ppt(meta_info, generation_content)
