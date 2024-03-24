import json

from pptx import Presentation


class PptxGenerator:

    @staticmethod
    def generate_ppt_content_prompt(topic, page_num):
        """生成PPT内容"""
        output_format = json.dumps({
            "title": "example title",
            "pages": [
                {
                    "title": "title for page 1",
                    "content": [
                        {
                            "title": "title for paragraph 1",
                            "description": "detail for paragraph 1",
                        },
                        {
                            "title": "title for paragraph 2",
                            "description": "detail for paragraph 2",
                        },
                    ],
                },
                {
                    "title": "title for page 2",
                    "content": [
                        {
                            "title": "title for paragraph 1",
                            "description": "detail for paragraph 1",
                        },
                        {
                            "title": "title for paragraph 2",
                            "description": "detail for paragraph 2",
                        },
                        {
                            "title": "title for paragraph 3",
                            "description": "detail for paragraph 3",
                        },
                    ],
                },
            ],
        }, ensure_ascii=True)

        return f'''我要准备1个关于{topic}的PPT，要求一共写{page_num}页，请你根据主题生成详细内容，不要省略。
        按这个JSON格式输出{output_format}，只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    @staticmethod
    def generate_ppt_file(ppt_content, author, save_path):
        """生成PPT文件"""
        ppt = Presentation()

        # PPT首页
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.placeholders[0].text = ppt_content['title']
        slide.placeholders[1].text = author

        # 内容页
        print(f'总共{len(ppt_content["pages"])}页...')
        for i, page in enumerate(ppt_content['pages']):
            print(f'生成第{i + 1}页:{page["title"]}')
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            # 标题
            slide.placeholders[0].text = page['title']
            # 正文
            for sub_content in page['content']:
                print(sub_content)
                # 一级正文
                sub_title = slide.placeholders[1].text_frame.add_paragraph()
                sub_title.text, sub_title.level = sub_content['title'], 1
                # 二级正文
                sub_description = slide.placeholders[1].text_frame.add_paragraph()
                sub_description.text, sub_description.level = sub_content['description'], 2

        ppt.save(save_path)
